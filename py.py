from pptx import Presentation
from pptx.util import Pt

# Create a new PowerPoint presentation
prs = Presentation()

# Function to add a slide with title and bullet points
def add_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    content.clear()  # Clear placeholder text
    for point in bullet_points:
        p = content.add_paragraph()
        p.text = point
        p.font.size = Pt(20)

# Slide 1: Why DevOps?
add_slide("Why Do We Need DevOps?", [
    "Faster software delivery",
    "Better collaboration between Development and Operations",
    "Increased deployment frequency",
    "Improved product quality and stability",
    "Faster issue resolution"
])

# Slide 2: Advantages of DevOps
add_slide("Advantages of DevOps", [
    "🚀 Faster time-to-market with automation",
    "🔄 Continuous Integration & Deployment (CI/CD)",
    "🔍 Early bug detection through continuous testing",
    "📊 Better monitoring and performance",
    "🤝 Cross-team collaboration and ownership"
])

# Slide 3: Disadvantages of DevOps
add_slide("Challenges / Disadvantages of DevOps", [
    "📚 Requires cultural change and training",
    "⚙️ Complex toolchain management",
    "💸 Higher initial setup cost",
    "🛡️ Security risks if not integrated properly",
    "🤯 Not suitable for all teams or projects"
])

# Save the presentation to an existing folder on your PC
prs.save(r"C:\Users\konderu Nagasai\sai\Why_DevOps_Presentation.pptx")

print("Presentation saved successfully!")
